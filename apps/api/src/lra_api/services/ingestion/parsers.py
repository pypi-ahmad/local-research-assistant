"""Document parsing functions by source type."""

from __future__ import annotations

import csv
import io
from pathlib import Path

import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pypdf import PdfReader
from pptx import Presentation


class ParserError(RuntimeError):
    """Raised when source parsing fails."""


def parse_bytes(content: bytes, filename: str, mime_type: str) -> str:
    """Parse source bytes into plain text."""
    suffix = Path(filename).suffix.lower()

    # Image/scanned sources are handled by OCR fallback in ingestion pipeline.
    if suffix in {".png", ".jpg", ".jpeg", ".webp", ".bmp"} or mime_type.startswith("image/"):
        return ""

    if suffix in {".csv"}:
        return _parse_csv(content)

    if suffix in {".txt", ".md"} or (mime_type.startswith("text/") and suffix not in {".csv"}):
        return content.decode("utf-8", errors="ignore")

    if suffix == ".pdf":
        return _parse_pdf(content)

    if suffix == ".docx":
        return _parse_docx(content)

    if suffix == ".pptx":
        return _parse_pptx(content)

    if suffix in {".xlsx", ".xls"}:
        return _parse_excel(content)

    if suffix in {".html", ".htm"}:
        return _parse_html(content)

    try:
        from docling.document_converter import DocumentConverter

        converter = DocumentConverter()
        result = converter.convert(io.BytesIO(content), source=filename)
        return result.document.export_to_markdown()
    except Exception as exc:  # noqa: BLE001
        raise ParserError(f"Unsupported or unreadable file: {filename}") from exc


def _parse_pdf(content: bytes) -> str:
    reader = PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages)


def _parse_docx(content: bytes) -> str:
    doc = DocxDocument(io.BytesIO(content))
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)


def _parse_pptx(content: bytes) -> str:
    presentation = Presentation(io.BytesIO(content))
    lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)


def _parse_csv(content: bytes) -> str:
    stream = io.StringIO(content.decode("utf-8", errors="ignore"))
    reader = csv.reader(stream)
    rows = [", ".join(row) for row in reader]
    return "\n".join(rows)


def _parse_excel(content: bytes) -> str:
    frame = pd.read_excel(io.BytesIO(content))
    return frame.to_csv(index=False)


def _parse_html(content: bytes) -> str:
    soup = BeautifulSoup(content.decode("utf-8", errors="ignore"), "html.parser")
    return soup.get_text(separator="\n")
